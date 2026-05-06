import json
import pickle
import csv
import yaml
import xml.etree.ElementTree as ET
import os
import mimetypes
import re
import time
from typing import Dict, Any, List, Union, Optional
from pathlib import Path
from abc import ABC, abstractmethod
from urllib.parse import parse_qs, unquote, urlparse

import requests

# For handling various file types
try:
    import pymupdf
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from PIL import Image
    PILLOW_AVAILABLE = True
except ImportError:
    PILLOW_AVAILABLE = False

try:
    from openpyxl import Workbook, load_workbook
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation
    # from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from ..core.module import BaseModule
from ..core.logging import logger

PREFIX_MATCH_THRESHOLD = 0.8
MIN_BASE_LEN_FOR_CHECK = 20

class StorageBase(BaseModule, ABC):
    """
    Abstract base class for comprehensive storage operations supporting various file types.
    Provides unified interface for local and remote storage operations.
    """
    
    def __init__(self, base_path: str = ".", **kwargs):
        """
        Initialize the StorageBase with configuration options.
        
        Args:
            base_path (str): Base directory for storage operations (default: current directory)
            **kwargs: Additional keyword arguments for parent class initialization
        """
        super().__init__(**kwargs)
        self.base_path = base_path
        self.return_file_url = kwargs.get("return_file_url", True)
        
        # File types that support append operations
        self.appendable_formats = {
            '.txt': self._append_text,
            '.json': self._append_json,
            '.csv': self._append_csv,
            '.yaml': self._append_yaml,
            '.yml': self._append_yaml,
            '.pickle': self._append_pickle,
            '.xlsx': self._append_excel,
            '.docx': self._append_docx,
            '.pptx': self._append_pptx
        }
        
        # Initialize storage-specific setup
        self._initialize_storage()
    
    @abstractmethod
    def _initialize_storage(self):
        """
        Initialize storage-specific setup. Override in subclasses for storage-specific initialization.
        """
        pass
    
    
    # Abstract methods that must be implemented by subclasses
    @abstractmethod
    def _read_raw(self, path: str, **kwargs) -> bytes:
        """Read raw file content - must be implemented by subclasses"""
        pass
    
    @abstractmethod
    def _write_raw(self, path: str, content: bytes, **kwargs) -> bool:
        """Write raw file content - must be implemented by subclasses"""
        pass
    
    @abstractmethod
    def _delete_raw(self, path: str) -> bool:
        """Delete file or directory - must be implemented by subclasses"""
        pass
    
    @abstractmethod
    def _list_raw(self, path: str = None, **kwargs) -> List[Dict[str, Any]]:
        """List files and directories - must be implemented by subclasses"""
        pass
    
    @abstractmethod
    def _exists_raw(self, path: str) -> bool:
        """Check if path exists - must be implemented by subclasses"""
        pass
    
    @abstractmethod
    def _create_directory_raw(self, path: str) -> bool:
        """Create directory - must be implemented by subclasses"""
        pass
    
    
    # ____________________ PATH TRANSLATION ____________________ #
    def translate_in(self, file_path: str) -> str:
        """
        Translate input file path by combining it with base_path.
        This method takes a user-provided path and converts it to the full system path.
        
        Args:
            file_path (str): User-provided file path (can be relative or absolute)
            
        Returns:
            str: Full system path combining base_path and file_path
        """
        # If the path is already absolute, return as is
        if os.path.isabs(file_path):
            return file_path
        
        # Always combine base_path with file_path to ensure working directory is respected
        # Check if this is a remote storage handler (like Supabase)
        if hasattr(self, 'bucket_name') and hasattr(self, 'supabase'):
            # For remote storage, treat base_path as a prefix within the bucket.
            # Normalize both base and incoming path to avoid double-prefixing and leading slashes.
            if self.base_path.startswith('/'):
                clean_base = self.base_path.lstrip('/')
            else:
                clean_base = self.base_path

            incoming = file_path.lstrip('/')

            # If the incoming path already contains the base prefix, avoid adding it again.
            if clean_base:
                if incoming == clean_base or incoming.startswith(f"{clean_base}/"):
                    return incoming
                else:
                    return f"{clean_base}/{incoming}"
            else:
                # No base prefix configured (e.g., base_path is '/'), just return normalized incoming
                return incoming
        else:
            # For local storage, use os.path.join for proper filesystem handling
            combined_path = os.path.join(self.base_path, file_path)
            normalized_path = os.path.normpath(combined_path)
            return normalized_path
    
    def translate_out(self, full_path: str) -> str:
        """
        Translate output full path by removing the base_path prefix.
        This method takes a full system path and converts it back to the user-relative path.
        
        Args:
            full_path (str): Full system path
            
        Returns:
            str: User-relative path with base_path removed
        """
        # If base_path is just "." or empty, return the full_path as is
        if self.base_path in [".", "", None]:
            return full_path
        
        # Check if this is a remote storage handler (like Supabase)
        if hasattr(self, 'bucket_name') and hasattr(self, 'supabase'):
            # For remote storage, handle path prefix removal robustly.
            # Normalize both base and full_path to avoid issues with leading '/'.
            if self.base_path.startswith('/'):
                clean_base = self.base_path.lstrip('/')
            else:
                clean_base = self.base_path

            incoming = full_path.lstrip('/')

            if clean_base:
                if incoming.startswith(f"{clean_base}/"):
                    relative_path = incoming[len(f"{clean_base}/"):]
                    return relative_path
                elif incoming == clean_base:
                    return ""
                else:
                    if len(clean_base) >= MIN_BASE_LEN_FOR_CHECK:
                        t = int(len(clean_base) * PREFIX_MATCH_THRESHOLD)
                        if incoming[:t] == clean_base[:t]:
                            return os.path.relpath(incoming, start=clean_base)
                        else:
                            return incoming
                    return incoming
            else:
                # No base prefix configured (e.g., base_path is '/'); return path without leading slash
                return incoming
        else:
            # For local storage, use os.path operations for proper filesystem handling
            # Convert both paths to absolute paths for comparison
            base_abs = os.path.abspath(self.base_path)
            full_abs = os.path.abspath(full_path)
            
            # Check if the full_path starts with base_path
            if full_abs.startswith(base_abs):
                relative_path = full_abs[len(base_abs):]
                if relative_path.startswith(os.sep):
                    relative_path = relative_path[1:]
                return relative_path
            
            t = int(len(base_abs) * MIN_BASE_LEN_FOR_CHECK / MIN_BASE_LEN_FOR_CHECK * PREFIX_MATCH_THRESHOLD) if len(base_abs) >= MIN_BASE_LEN_FOR_CHECK else 0
            if t and full_abs[:t] == base_abs[:t]:
                return os.path.relpath(full_abs, start=base_abs)
            else:
                return full_path
    
    
    # ____________________ FILE INFO ____________________ #
    def get_file_type(self, file_path: str) -> str:
        """Get the file extension from a file path"""
        return Path(file_path).suffix.lower()

    def get_file_info(self, file_path: str) -> Dict[str, Any]:
        """Get comprehensive information about a file"""
        try:
            target_path = self.translate_in(file_path)
            if not self._exists_raw(target_path):
                return {"success": False, "error": f"File {file_path} does not exist"}
            
            # For now, return basic info - subclasses can override for more details
            res = {
                "success": True,
                "file_path": target_path,
                "file_name": Path(target_path).name,
                "file_extension": Path(target_path).suffix.lower(),
                "exists": True
            }
            if self.return_file_url:
                try:
                    res["url"] = self._get_file_url(target_path)
                except Exception:
                    pass
            return res
        except Exception as e:
            logger.error(f"Error getting file info for {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    def create_directory(self, path: str) -> Dict[str, Any]:
        """Create directory"""
        try:
            target_path = self.translate_in(path)
            success = self._create_directory_raw(target_path)
            if success:
                res = {"success": True, "path": target_path, "message": "Directory created successfully"}
                if self.return_file_url:
                    try:
                        res["url"] = self._get_file_url(target_path)
                    except Exception:
                        pass
                return res
            else:
                return {"success": False, "error": "Failed to create directory", "path": target_path}
        except Exception as e:
            logger.error(f"Error creating directory {path}: {str(e)}")
            return {"success": False, "error": str(e), "path": path}
    
    def exists(self, path: str) -> bool:
        """Check if path exists"""
        target_path = self.translate_in(path)
        return self._exists_raw(target_path)

    def _get_unique_name(self, translated_path: str) -> str:
        """
        Generate a unique file path by appending an incrementing suffix
        before the extension when a file with the same name already exists.

        Input contract:
        - `translated_path` MUST already be translated via `translate_in`.
        - Works for both local and remote handlers; existence is checked
          using `_exists_raw` against the provided path.

        Examples:
        - "/base/dir/test.py" -> "/base/dir/test_1.py" (if exists)
        - "/base/dir/test" -> "/base/dir/test_1" (no extension case)
        """
        try:
            p = Path(translated_path)
            parent = p.parent
            stem = p.stem
            suffix = p.suffix

            candidate = p
            counter = 1
            while self._exists_raw(str(candidate)):
                candidate = parent / f"{stem}_{counter}{suffix}"
                counter += 1

            return str(candidate)
        except Exception:
            # If anything goes wrong, fall back to original path
            return translated_path
    
    
    # ____________________ CURD ____________________ #
    def delete(self, path: str) -> Dict[str, Any]:
        """Delete file or directory"""
        try:
            target_path = self.translate_in(path)
            success = self._delete_raw(target_path)
            if success:
                res = {"success": True, "path": target_path, "message": "Deleted successfully"}
                if self.return_file_url:
                    try:
                        res["url"] = self._get_file_url(target_path)
                    except Exception:
                        pass
                return res
            else:
                return {"success": False, "error": "Failed to delete", "path": target_path}
        except Exception as e:
            logger.error(f"Error deleting {path}: {str(e)}")
            return {"success": False, "error": str(e), "path": path}
    
    def move(self, source: str, destination: str) -> Dict[str, Any]:
        """Move/rename file or directory"""
        try:
            resolved_source = self.translate_in(source)
            # Ensure destination uniqueness unconditionally (destination is translated)
            desired_destination = self.translate_in(destination)
            resolved_destination = self._get_unique_name(desired_destination)
            
            # Read source content
            content = self._read_raw(resolved_source)
            
            # Write to destination
            success = self._write_raw(resolved_destination, content)
            if success:
                # Delete source
                self._delete_raw(resolved_source)
                # Return only destination info, using usual keys
                res = {"success": True, "file_path": resolved_destination, "message": "Moved successfully"}
                if self.return_file_url:
                    try:
                        res["url"] = self._get_file_url(resolved_destination)
                    except Exception:
                        pass
                return res
            else:
                return {"success": False, "error": "Failed to write to destination", "source": resolved_source, "destination": resolved_destination}
        except Exception as e:
            logger.error(f"Error moving {source} to {destination}: {str(e)}")
            return {"success": False, "error": str(e), "source": source, "destination": destination}
    
    def copy(self, source: str, destination: str) -> Dict[str, Any]:
        """Copy file"""
        try:
            resolved_source = self.translate_in(source)
            # Ensure destination uniqueness unconditionally (destination is translated)
            desired_destination = self.translate_in(destination)
            resolved_destination = self._get_unique_name(desired_destination)
            
            # Read source content
            content = self._read_raw(resolved_source)
            
            # Write to destination
            success = self._write_raw(resolved_destination, content)
            if success:
                # Return only destination info, using usual keys
                res = {"success": True, "file_path": resolved_destination, "message": "Copied successfully"}
                if self.return_file_url:
                    try:
                        res["url"] = self._get_file_url(resolved_destination)
                    except Exception:
                        pass
                return res
            else:
                return {"success": False, "error": "Failed to write to destination", "source": resolved_source, "destination": resolved_destination}
        except Exception as e:
            logger.error(f"Error copying {source} to {destination}: {str(e)}")
            return {"success": False, "error": str(e), "source": source, "destination": destination}
    
    def list(self, path: str = None, max_depth: int = 3, include_hidden: bool = False) -> Dict[str, Any]:
        """List files and directories"""
        try:
            target_path = self.translate_in(path) if path else str(self.base_path)
            items = self._list_raw(target_path, max_depth=max_depth, include_hidden=include_hidden)
            if self.return_file_url:
                for item in items:
                    try:
                        item_path = item.get("path")
                        if item_path:
                            item["url"] = self._get_file_url(item_path)
                    except Exception:
                        continue
            
            output = []
            for item in items:
                item["path"] = self.translate_out(item.get("path"))
                output.append(item)
            return {
                "success": True,
                "path": target_path,
                "items": output,
                "total_count": len(items)
            }
        except Exception as e:
            logger.error(f"Error listing {path}: {str(e)}")
            return {"success": False, "error": str(e), "path": path}
    
    def save(self, file_path: str, content: Any, **kwargs) -> Dict[str, Any]:
        """
        Save content to a file with automatic format detection.
        This method replaces the old save method with the improved create_file logic.
        
        Args:
            file_path (str): Path where the file should be saved
            content (Any): Content to save to the file
            **kwargs: Additional arguments for file creation (encoding, format, etc.)
            
        Returns:
            Dict[str, Any]: Result of the operation with success status and details
        """
        try:
            # Get file type to determine the appropriate save method
            file_extension = self.get_file_type(file_path)
            # Always ensure unique, path-aware destination using translated path
            desired_path = self.translate_in(file_path)
            target_file_path = self._get_unique_name(desired_path)

            # If content is already bytes, write it raw regardless of extension
            if isinstance(content, bytes):
                success = self._write_raw(target_file_path, content, **kwargs)
                if success:
                    res = {
                        "success": True,
                        "message": "File saved successfully",
                        "file_path": self.translate_out(target_file_path),
                        "full_path": target_file_path,
                        "size": len(content)
                    }
                    if self.return_file_url:
                        res["url"] = self._get_file_url(target_file_path)
                    return res
                else:
                    return {
                        "success": False,
                        "message": f"Failed to save file '{file_path}'",
                        "file_path": self.translate_out(file_path),
                        "full_path": target_file_path
                    }
            
            # Route to specialized save methods based on file type
            if file_extension == '.json':
                return self._save_json(target_file_path, content, **kwargs)
            elif file_extension in ['.txt', '.md', '.log']:
                return self._save_text(target_file_path, content, **kwargs)
            elif file_extension == '.csv':
                return self._save_csv(target_file_path, content, **kwargs)
            elif file_extension in ['.yaml', '.yml']:
                return self._save_yaml(target_file_path, content, **kwargs)
            elif file_extension == '.xml':
                return self._save_xml(target_file_path, content, **kwargs)
            elif file_extension == '.xlsx':
                return self._save_excel(target_file_path, content, **kwargs)
            elif file_extension == '.docx':
                return self._save_docx(target_file_path, content, **kwargs)
            elif file_extension == '.pptx':
                return self._save_pptx(target_file_path, content, **kwargs)
            elif file_extension == '.pickle':
                return self._save_pickle(target_file_path, content, **kwargs)
            elif file_extension == '.pdf':
                return self._save_pdf(target_file_path, content, **kwargs)
            elif file_extension in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff']:
                return self._save_image(target_file_path, content, **kwargs)
            else:
                # For other file types, use the generic approach
                # Convert content to bytes if it's not already
                if isinstance(content, str):
                    content_bytes = content.encode(kwargs.get('encoding', 'utf-8'))
                elif isinstance(content, bytes):
                    content_bytes = content
                else:
                    content_bytes = str(content).encode(kwargs.get('encoding', 'utf-8'))
                
                # Write the file using the raw method
                success = self._write_raw(target_file_path, content_bytes, **kwargs)
            
                if success:
                    res = {
                        "success": True,
                        "message": "File saved successfully",
                        "file_path": self.translate_out(target_file_path),
                        "full_path": target_file_path,
                        "size": len(content_bytes)
                    }
                    if self.return_file_url:
                        res["url"] = self._get_file_url(target_file_path)
                    return res
                else:
                    return {
                        "success": False,
                        "message": f"Failed to save file '{file_path}'",
                        "file_path": self.translate_out(file_path),
                        "full_path": target_file_path
                    }
        except Exception as e:
            logger.error(f"Error saving file {file_path}: {str(e)}")
            return {
                "success": False,
                "message": f"Error saving file: {str(e)}",
                "file_path": file_path
            }
    
    def read(self, file_path: Optional[str] = None, url: Optional[str] = None, **kwargs) -> Dict[str, Any]:
        """
        Read content from a file with automatic format detection.

        Supports two modes of operation:
        1. File path mode: Read from storage (local/remote)
        2. URL mode: Download from URL to memory and read directly (no storage I/O)

        Args:
            file_path (Optional[str]): Path to file in storage (relative to base_path)
            url (Optional[str]): URL to download file from (downloads to memory only)
            **kwargs: Additional arguments passed to format-specific readers

        Returns:
            Dict[str, Any]: Result dictionary with structure:
                - success (bool): Whether the operation succeeded
                - content (Any): Parsed content (type depends on file format)
                - file_path (str): File path used (provided or extracted from URL)
                - url (str, optional): Original URL (only when url parameter is used)
                - Additional format-specific fields (e.g., rows, sheet_name, etc.)

        Examples:
            >>> # Example 1: Read from storage
            >>> result = storage.read(file_path="data/config.json")
            >>> # Returns: {
            >>> #   "success": True,
            >>> #   "content": {"key": "value"},
            >>> #   "file_path": "data/config.json"
            >>> # }

            >>> # Example 2: Read from URL (downloads to memory, no storage I/O)
            >>> result = storage.read(url="https://example.com/data.json")
            >>> # Returns: {
            >>> #   "success": True,
            >>> #   "content": {"key": "value"},
            >>> #   "file_path": "data.json",
            >>> #   "url": "https://example.com/data.json"
            >>> # }

        Note:
            When using URL mode, the file is downloaded to memory and read directly
            without being saved to storage. This is more efficient for temporary reads.
        """
        if not file_path and not url:
            return {
                "success": False,
                "message": "Input must be provided"
            }
        
        try:
            url_download_used = False
            if url and not file_path:
                print("Downloading file from URL...")
                url_download_used = True

                download_result = self._download_content(url, timeout=30, max_retries=3, memory_only=True)
                if not download_result.get('success', False):
                    return {
                        "success": False,
                        "error": download_result.get('error', 'Download failed')
                    }

                if 'content' in download_result:
                    kwargs['content'] = download_result['content']
                    file_path = download_result['filename']
                else:
                    file_path = download_result.get('file_path', None)
            
            if not file_path:
                return {
                    "success": False,
                    "message": "File path must be provided"
                }
            
            target_file_path = self.translate_in(file_path)
            file_extension = Path(target_file_path).suffix.lower()
            
            # Handle different file types
            
            print(f"While reading the file: {target_file_path}, file extension: {file_extension}")
            
            if file_extension == '.json':
                result = self._read_json(target_file_path, **kwargs)
            elif file_extension in ['.yaml', '.yml']:
                result = self._read_yaml(target_file_path, **kwargs)
            elif file_extension == '.csv':
                result = self._read_csv(target_file_path, **kwargs)
            elif file_extension == '.xlsx':
                result = self._read_excel(target_file_path, **kwargs)
            elif file_extension == '.xml':
                result = self._read_xml(target_file_path, **kwargs)
            elif file_extension == '.docx':
                result = self._read_docx(target_file_path, **kwargs)
            elif file_extension == '.pptx':
                result = self._read_pptx(target_file_path, **kwargs)
            elif file_extension == '.pickle':
                result = self._read_pickle(target_file_path, **kwargs)
            elif file_extension == '.pdf':
                result = self._read_pdf(target_file_path, **kwargs)
            elif file_extension in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff']:
                result = self._read_image(target_file_path, **kwargs)
            else:
                # Default to text
                print(f"Defaulting to text read for file: {target_file_path}")
                result = self._read_text(target_file_path, **kwargs)

            if url_download_used and isinstance(result, dict):
                result["url"] = url

            return result
            
        except Exception as e:
            logger.error(f"Error reading {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    def append(self, file_path: str, content: Any, **kwargs) -> Dict[str, Any]:
        """Append content to a file (only for supported formats)"""
        try:
            target_file_path = self.translate_in(file_path)
            file_extension = Path(target_file_path).suffix.lower()
            
            if file_extension in self.appendable_formats:
                return self.appendable_formats[file_extension](target_file_path, content, **kwargs)
            else:
                return {"success": False, "error": f"Append not supported for {file_extension} files"}
                
        except Exception as e:
            logger.error(f"Error appending to {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    # Text file handlers
    def _save_text(self, file_path: str, content: Any, encoding: str = 'utf-8', **kwargs) -> Dict[str, Any]:
        """Save text content to a file"""
        try:
            # Convert content to bytes
            if isinstance(content, str):
                content_bytes = content.encode(encoding)
            else:
                content_bytes = str(content).encode(encoding)
            
            # Use raw write method
            success = self._write_raw(file_path, content_bytes, **kwargs)
            
            if success:
                res = {
                    "success": True,
                    "message": "File saved successfully",
                    "file_path": file_path,
                    "content_length": len(content_bytes)
                }
                if self.return_file_url:
                    res["url"] = self._get_file_url(file_path)
                return res
            else:
                return {"success": False, "error": "Failed to write file", "file_path": file_path}
        except Exception as e:
            logger.error(f"Error saving text file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    def _read_text(self, file_path: Optional[str] = None, content: Any = None, encoding: str = 'utf-8', **kwargs) -> Dict[str, Any]:
        """Read text content from a file"""
        try:
            # Use raw read method
            if content is None:
                if not file_path:
                    return {"success": False, "error": "Either file_path or content must be provided"}
                content = self._read_raw(file_path, **kwargs)
            
            if not isinstance(content, bytes):
                return {"success": False, "error": "Content must be bytes"}
            
            content_bytes = content
            text_content = content_bytes.decode(encoding)
            
            base = {
                "success": True,
                "content": text_content,
                "file_path": file_path,
                "content_length": len(text_content)
            }
            return base
        except Exception as e:
            logger.error(f"Error reading text file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
        
    def _append_text(self, file_path: str, content: str, encoding: str = 'utf-8', **kwargs) -> Dict[str, Any]:
        """Append text content to a file"""
        try:
            # Convert content to bytes
            content_bytes = str(content).encode(encoding)
            
            # For append, read existing content first
            existing_bytes = b""
            if self._exists_raw(file_path):
                existing_bytes = self._read_raw(file_path, **kwargs)
            
            # Combine existing and new content
            combined_bytes = existing_bytes + content_bytes
            
            # Write combined content
            success = self._write_raw(file_path, combined_bytes, **kwargs)
            
            if success:
                res = {
                    "success": True,
                    "message": f"Content appended to file {file_path}",
                    "file_path": file_path
                }
                if self.return_file_url:
                    res["url"] = self._get_file_url(file_path)
                return res
            else:
                return {"success": False, "error": "Failed to append to file", "file_path": file_path}
        except Exception as e:
            logger.error(f"Error appending to text file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    # JSON file handlers
    def _save_json(self, file_path: str, content: Any, indent: int = 2, **kwargs) -> Dict[str, Any]:
        """Save JSON content to a file"""
        try:
            # Convert content to JSON string
            if isinstance(content, str):
                # Validate JSON
                json.loads(content)
                json_content = content
            else:
                json_content = json.dumps(content, indent=indent, ensure_ascii=False)
            
            # Convert to bytes
            content_bytes = json_content.encode('utf-8')
            
            # Use raw write method
            success = self._write_raw(file_path, content_bytes, **kwargs)
            
            if success:
                res = {
                    "success": True,
                    "message": "JSON file saved successfully",
                    "file_path": file_path
                }
                if self.return_file_url:
                    res["url"] = self._get_file_url(file_path)
                return res
            else:
                return {"success": False, "error": "Failed to write file", "file_path": file_path}
        except Exception as e:
            logger.error(f"Error saving JSON file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    def _read_json(self, file_path: Optional[str] = None, content: Any = None, **kwargs) -> Dict[str, Any]:
        """Read JSON content from a file"""
        try:
            # Use raw read method
            if content is None:
                if not file_path:
                    return {"success": False, "error": "Either file_path or content must be provided"}
                content = self._read_raw(file_path, **kwargs)
            
            if not isinstance(content, bytes):
                return {"success": False, "error": "Content must be bytes"}

            content_bytes = content
            content_str = content_bytes.decode('utf-8')
            
            # Parse JSON
            json_content = json.loads(content_str)
            
            base = {
                "success": True,
                "content": json_content,
                "file_path": file_path
            }
            return base
        except Exception as e:
            logger.error(f"Error reading JSON file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    def _append_json(self, file_path: str, content: Any, **kwargs) -> Dict[str, Any]:
        """Append content to JSON file (for arrays)"""
        try:
            # Read existing content
            existing_content = []
            if self._exists_raw(file_path):
                existing_bytes = self._read_raw(file_path, **kwargs)
                existing_str = existing_bytes.decode('utf-8')
                existing_content = json.loads(existing_str)
            
            # Merge content
            if isinstance(existing_content, list):
                if isinstance(content, list):
                    existing_content.extend(content)
                else:
                    existing_content.append(content)
            elif isinstance(existing_content, dict):
                if isinstance(content, dict):
                    existing_content.update(content)
                else:
                    return {"success": False, "error": "Cannot append non-dict to JSON dict"}
            else:
                existing_content = [existing_content]
                if isinstance(content, list):
                    existing_content.extend(content)
                else:
                    existing_content.append(content)
            
            # Convert to JSON string and bytes
            json_content = json.dumps(existing_content, indent=2, ensure_ascii=False)
            content_bytes = json_content.encode('utf-8')
            
            # Write combined content
            success = self._write_raw(file_path, content_bytes, **kwargs)
            
            if success:
                res = {
                    "success": True,
                    "message": f"Content appended to JSON file {file_path}",
                    "file_path": file_path
                }
                if self.return_file_url:
                    res["url"] = self._get_file_url(file_path)
                return res
            else:
                return {"success": False, "error": "Failed to append to file", "file_path": file_path}
        except Exception as e:
            logger.error(f"Error appending to JSON file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    # CSV file handlers
    def _save_csv(self, file_path: str, content: Any, **kwargs) -> Dict[str, Any]:
        """Save CSV content to a file - handles both raw CSV strings and structured data"""
        try:
            if not content:
                return {"success": False, "error": "No content to save"}
            
            from io import StringIO
            
            # Build CSV content in memory
            csv_buffer = StringIO()
            
            # If content is a string, use it directly
            if isinstance(content, str):
                csv_content = content
                rows = content.count('\n')
            # If content is a list of dictionaries, use CSV writer
            elif isinstance(content, list) and content and isinstance(content[0], dict):
                fieldnames = content[0].keys()
                writer = csv.DictWriter(csv_buffer, fieldnames=fieldnames)
                writer.writeheader()
                writer.writerows(content)
                csv_content = csv_buffer.getvalue()
                rows = len(content)
            # If content is a list of lists, use CSV writer
            elif isinstance(content, list) and content and isinstance(content[0], list):
                writer = csv.writer(csv_buffer)
                writer.writerows(content)
                csv_content = csv_buffer.getvalue()
                rows = len(content)
            else:
                return {"success": False, "error": "CSV content must be a string, list of dictionaries, or list of lists"}
            
            # Convert to bytes and write
            content_bytes = csv_content.encode('utf-8')
            success = self._write_raw(file_path, content_bytes, **kwargs)
            
            if success:
                res = {
                    "success": True,
                    "message": "CSV file saved successfully",
                    "file_path": file_path,
                    "rows": rows
                }
                if self.return_file_url:
                    res["url"] = self._get_file_url(file_path)
                return res
            else:
                return {"success": False, "error": "Failed to write file", "file_path": file_path}
            
        except Exception as e:
            logger.error(f"Error saving CSV file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    def _read_csv(self, file_path: Optional[str] = None, content: Any = None, **kwargs) -> Dict[str, Any]:
        """Read CSV content from a file"""
        try:
            # Use raw read method
            if content is None:
                if not file_path:
                    return {"success": False, "error": "Either file_path or content must be provided"}
                content = self._read_raw(file_path, **kwargs)
            
            if not isinstance(content, bytes):
                return {"success": False, "error": "Content must be bytes"}

            content_bytes = content
            content_str = content_bytes.decode('utf-8')
            
            # Parse CSV
            from io import StringIO
            reader = csv.DictReader(StringIO(content_str))
            csv_content = list(reader)
            
            base = {
                "success": True,
                "content": csv_content,
                "file_path": file_path,
                "rows": len(csv_content)
            }
            return base
        except Exception as e:
            logger.error(f"Error reading CSV file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    def _append_csv(self, file_path: str, content: List[Dict[str, Any]], **kwargs) -> Dict[str, Any]:
        """Append content to CSV file"""
        try:
            if not content:
                return {"success": False, "error": "No content to append"}
            
            # Read existing content
            existing_content = []
            if self._exists_raw(file_path):
                existing_bytes = self._read_raw(file_path, **kwargs)
                existing_str = existing_bytes.decode('utf-8')
                from io import StringIO
                reader = csv.DictReader(StringIO(existing_str))
                existing_content = list(reader)
            
            # Combine content
            combined_content = existing_content + content
            
            # Write combined content
            from io import StringIO
            csv_buffer = StringIO()
            if combined_content:
                fieldnames = combined_content[0].keys()
                writer = csv.DictWriter(csv_buffer, fieldnames=fieldnames)
                writer.writeheader()
                writer.writerows(combined_content)
            
            csv_content = csv_buffer.getvalue()
            content_bytes = csv_content.encode('utf-8')
            success = self._write_raw(file_path, content_bytes, **kwargs)
            
            if success:
                res = {
                    "success": True,
                    "message": f"Content appended to CSV file {file_path}",
                    "file_path": file_path,
                    "appended_rows": len(content)
                }
                if self.return_file_url:
                    res["url"] = self._get_file_url(file_path)
                return res
            else:
                return {"success": False, "error": "Failed to append to file", "file_path": file_path}
        except Exception as e:
            logger.error(f"Error appending to CSV file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    # YAML file handlers
    def _save_yaml(self, file_path: str, content: Any, **kwargs) -> Dict[str, Any]:
        """Save YAML content to a file"""
        try:
            # Convert content to YAML string
            yaml_content = yaml.dump(content, default_flow_style=False, allow_unicode=True)
            content_bytes = yaml_content.encode('utf-8')
            
            # Use raw write method
            success = self._write_raw(file_path, content_bytes, **kwargs)
            
            if success:
                res = {
                    "success": True,
                    "message": "YAML file saved successfully",
                    "file_path": file_path
                }
                if self.return_file_url:
                    res["url"] = self._get_file_url(file_path)
                return res
            else:
                return {"success": False, "error": "Failed to write file", "file_path": file_path}
        except Exception as e:
            logger.error(f"Error saving YAML file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    def _read_yaml(self, file_path: Optional[str] = None, content: Any = None, **kwargs) -> Dict[str, Any]:
        """Read YAML content from a file"""
        try:
            # Use raw read method
            if content is None:
                if not file_path:
                    return {"success": False, "error": "Either file_path or content must be provided"}
                content = self._read_raw(file_path, **kwargs)
            
            if not isinstance(content, bytes):
                return {"success": False, "error": "Content must be bytes"}

            content_bytes = content
            content_str = content_bytes.decode('utf-8')
            
            # Parse YAML
            yaml_content = yaml.safe_load(content_str)
            
            base = {
                "success": True,
                "content": yaml_content,
                "file_path": file_path
            }
            return base
        except Exception as e:
            logger.error(f"Error reading YAML file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    def _append_yaml(self, file_path: str, content: Any, **kwargs) -> Dict[str, Any]:
        """Append content to YAML file (for lists)"""
        try:
            # Read existing content
            existing_content = []
            if self._exists_raw(file_path):
                existing_bytes = self._read_raw(file_path, **kwargs)
                existing_str = existing_bytes.decode('utf-8')
                existing_content = yaml.safe_load(existing_str) or []
            
            # Merge content
            if isinstance(existing_content, list):
                if isinstance(content, list):
                    existing_content.extend(content)
                else:
                    existing_content.append(content)
            elif isinstance(existing_content, dict):
                if isinstance(content, dict):
                    existing_content.update(content)
                else:
                    return {"success": False, "error": "Cannot append non-dict to YAML dict"}
            else:
                existing_content = [existing_content]
                if isinstance(content, list):
                    existing_content.extend(content)
                else:
                    existing_content.append(content)
            
            # Convert to YAML string and bytes
            yaml_content = yaml.dump(existing_content, default_flow_style=False, allow_unicode=True)
            content_bytes = yaml_content.encode('utf-8')
            
            # Write combined content
            success = self._write_raw(file_path, content_bytes, **kwargs)
            
            if success:
                res = {
                    "success": True,
                    "message": f"Content appended to YAML file {file_path}",
                    "file_path": file_path
                }
                if self.return_file_url:
                    res["url"] = self._get_file_url(file_path)
                return res
            else:
                return {"success": False, "error": "Failed to append to file", "file_path": file_path}
        except Exception as e:
            logger.error(f"Error appending to YAML file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    # XML file handlers
    def _save_xml(self, file_path: str, content: Any, root_tag: str = "root", **kwargs) -> Dict[str, Any]:
        """Save XML content to a file"""
        try:
            # If content is already a string (raw XML), use it directly
            if isinstance(content, str):
                # Check if it's already valid XML
                try:
                    ET.fromstring(content)
                    xml_content = content
                except ET.ParseError:
                    # Not valid XML, treat as text content and wrap it
                    root = ET.Element(root_tag)
                    root.text = content
                    xml_content = ET.tostring(root, encoding='unicode')
            # If content is a dictionary, convert to XML
            elif isinstance(content, dict):
                def dict_to_xml(data, root):
                    for key, value in data.items():
                        child = ET.SubElement(root, key)
                        if isinstance(value, dict):
                            dict_to_xml(value, child)
                        else:
                            child.text = str(value)
                
                root = ET.Element(root_tag)
                dict_to_xml(content, root)
                xml_content = ET.tostring(root, encoding='unicode')
            else:
                # For other types, wrap in root element
                root = ET.Element(root_tag)
                root.text = str(content)
                xml_content = ET.tostring(root, encoding='unicode')
            
            # Convert to bytes and write
            content_bytes = xml_content.encode('utf-8')
            success = self._write_raw(file_path, content_bytes, **kwargs)
            
            if success:
                res = {
                    "success": True,
                    "message": "XML file saved successfully",
                    "file_path": file_path
                }
                if self.return_file_url:
                    res["url"] = self._get_file_url(file_path)
                return res
            else:
                return {"success": False, "error": "Failed to write file", "file_path": file_path}
        except Exception as e:
            logger.error(f"Error saving XML file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    def _read_xml(self, file_path: Optional[str] = None, content: Any = None, **kwargs) -> Dict[str, Any]:
        """Read XML content from a file"""
        try:
            # Use raw read method
            if content is None:
                if not file_path:
                    return {"success": False, "error": "Either file_path or content must be provided"}
                content = self._read_raw(file_path, **kwargs)
            
            if not isinstance(content, bytes):
                return {"success": False, "error": "Content must be bytes"}

            content_bytes = content
            content_str = content_bytes.decode('utf-8')
            
            # Parse XML
            root = ET.fromstring(content_str)
            
            def xml_to_dict(element):
                result = {}
                for child in element:
                    if len(child) == 0:
                        result[child.tag] = child.text
                    else:
                        result[child.tag] = xml_to_dict(child)
                return result
            
            xml_content = xml_to_dict(root)
            
            base = {
                "success": True,
                "content": xml_content,
                "file_path": file_path
            }
            return base
        except Exception as e:
            logger.error(f"Error reading XML file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    # Excel file handlers
    def _save_excel(self, file_path: str, content: List[List[Any]], sheet_name: str = "Sheet1", **kwargs) -> Dict[str, Any]:
        """Save Excel content to a file"""
        if not EXCEL_AVAILABLE:
            return {"success": False, "error": "openpyxl library not available"}
        
        try:
            from io import BytesIO
            
            # Create workbook in memory
            workbook = Workbook()
            worksheet = workbook.active
            worksheet.title = sheet_name
            
            for row in content:
                worksheet.append(row)
            
            # Save to BytesIO buffer
            buffer = BytesIO()
            workbook.save(buffer)
            content_bytes = buffer.getvalue()
            
            # Use raw write method
            success = self._write_raw(file_path, content_bytes, **kwargs)
            
            if success:
                res = {
                    "success": True,
                    "message": "Excel file saved successfully",
                    "file_path": file_path,
                    "rows": len(content)
                }
                if self.return_file_url:
                    res["url"] = self._get_file_url(file_path)
                return res
            else:
                return {"success": False, "error": "Failed to write file", "file_path": file_path}
        except Exception as e:
            logger.error(f"Error saving Excel file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    def _read_excel(self, file_path: Optional[str] = None, content: Any = None, sheet_name: str = None, **kwargs) -> Dict[str, Any]:
        """Read Excel content from a file"""
        if not EXCEL_AVAILABLE:
            return {"success": False, "error": "openpyxl library not available"}
        
        try:
            from io import BytesIO
            
            # Use raw read method
            if content is None:
                if not file_path:
                    return {"success": False, "error": "Either file_path or content must be provided"}
                content = self._read_raw(file_path, **kwargs)
            
            if not isinstance(content, bytes):
                return {"success": False, "error": "Content must be bytes"}

            content_bytes = content
            
            # Load workbook from bytes
            workbook = load_workbook(BytesIO(content_bytes), data_only=True)
            sheet_names = workbook.sheetnames
            
            if sheet_name is None:
                sheet_name = sheet_names[0]
            
            if sheet_name not in sheet_names:
                return {"success": False, "error": f"Sheet '{sheet_name}' not found"}
            
            worksheet = workbook[sheet_name]
            content = []
            
            def safe_convert(cell):
                if isinstance(cell, (str, int, float, bool, type(None))):
                    return cell
                try:
                    return str(cell)
                except Exception:
                    return None

            for row in worksheet.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    # Coerce all non-JSON-primitive values to str to ensure serializability
                    processed_row = [
                        safe_convert(cell)
                        for cell in row
                    ]
                    content.append(processed_row)
            
            base = {
                "success": True,
                "content": content,
                "file_path": file_path,
                "sheet_name": sheet_name,
                "rows": len(content)
            }
            return base
        except Exception as e:
            logger.error(f"Error reading Excel file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    def _append_excel(self, file_path: str, content: List[List[Any]], sheet_name: str = None, **kwargs) -> Dict[str, Any]:
        """Append content to Excel file"""
        if not EXCEL_AVAILABLE:
            return {"success": False, "error": "openpyxl library not available"}
        
        try:
            from io import BytesIO
            
            if not self._exists_raw(file_path):
                return self._save_excel(file_path, content, sheet_name or "Sheet1", **kwargs)
            
            # Read existing content
            content_bytes = self._read_raw(file_path, **kwargs)
            workbook = load_workbook(BytesIO(content_bytes))
            sheet_names = workbook.sheetnames
            
            if sheet_name is None:
                sheet_name = sheet_names[0]
            
            if sheet_name not in sheet_names:
                return {"success": False, "error": f"Sheet '{sheet_name}' not found"}
            
            worksheet = workbook[sheet_name]
            
            for row in content:
                worksheet.append(row)
            
            # Save to BytesIO buffer
            buffer = BytesIO()
            workbook.save(buffer)
            updated_bytes = buffer.getvalue()
            
            # Write updated content
            success = self._write_raw(file_path, updated_bytes, **kwargs)
            
            if success:
                res = {
                    "success": True,
                    "message": f"Content appended to Excel file {file_path}",
                    "file_path": file_path,
                    "appended_rows": len(content)
                }
                if self.return_file_url:
                    res["url"] = self._get_file_url(file_path)
                return res
            else:
                return {"success": False, "error": "Failed to append to file", "file_path": file_path}
        except Exception as e:
            logger.error(f"Error appending to Excel file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    
    def _read_docx(self, file_path: Optional[str] = None, content: Any = None, **kwargs) -> Dict[str, Any]:
        """Read content from a DOCX file using python-docx built-ins.

        Preserves simple structure without complex parsing:
        - paragraphs: flat list of paragraph texts (backward-compatible)
        - headings: list of {level, text} for Heading-styled paragraphs
        - tables: list of tables, each table is list of rows of cell texts
        - lists: best-effort list detection via style names containing 'List'
        """
        try:
            from io import BytesIO
            from docx import Document

            # Read raw bytes and load DOCX
            if content is None:
                if not file_path:
                    return {"success": False, "error": "Either file_path or content must be provided"}
                content = self._read_raw(file_path, **kwargs)
            
            if not isinstance(content, bytes):
                return {"success": False, "error": "Content must be bytes"}

            content_bytes = content
            document = Document(BytesIO(content_bytes))

            # Paragraphs and simple heading/list metadata
            paragraphs_text = []
            headings = []
            lists = []
            for para in document.paragraphs:
                text = para.text
                paragraphs_text.append(text)
                try:
                    style_name = para.style.name if para.style else None
                except Exception:
                    style_name = None
                if style_name and style_name.startswith("Heading") and text:
                    headings.append({"level": style_name, "text": text})
                if style_name and "List" in style_name and text:
                    lists.append(text)

            # Tables as plain matrices of cell text
            tables = []
            for table in document.tables:
                table_rows = []
                for row in table.rows:
                    table_rows.append([cell.text for cell in row.cells])
                tables.append(table_rows)

            base = {
                "success": True,
                "content": paragraphs_text,
                "file_path": file_path,
                "rows": len(paragraphs_text),
                "headings": headings,
                "tables": tables,
                "lists": lists,
            }
            return base
        except Exception as e:
            logger.error(f"Error reading DOCX file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    def _save_docx(self, file_path: str, content: Any, **kwargs) -> Dict[str, Any]:
        """
        Save content to a DOCX file.
        
        Input: String → paragraph, List → paragraphs/bullet points
        Output: Creates new DOCX document
        """
        try:
            from docx import Document
            from io import BytesIO
            
            # Create a new document
            doc = Document()
            
            # Handle different content types
            if isinstance(content, str):
                # Single paragraph
                doc.add_paragraph(content)
            
            elif isinstance(content, list):
                # List of paragraphs or lists
                for item in content:
                    if isinstance(item, str):
                        # Add as plain text paragraph
                        doc.add_paragraph(item)
                    elif isinstance(item, list):
                        # Add as bullet list
                        for list_item in item:
                            if isinstance(list_item, str):
                                # Try to use bullet style, fall back to normal paragraph with bullet
                                try:
                                    doc.add_paragraph(list_item, style='List Bullet')
                                except Exception:
                                    doc.add_paragraph(f"• {list_item}")
            
            # Save document to bytes
            doc_buffer = BytesIO()
            doc.save(doc_buffer)
            doc_bytes = doc_buffer.getvalue()
            
            # Use raw write method
            success = self._write_raw(file_path, doc_bytes, **kwargs)
            
            if success:
                res = {
                    "success": True,
                    "message": "DOCX file saved successfully",
                    "file_path": file_path
                }
                if self.return_file_url:
                    res["url"] = self._get_file_url(file_path)
                return res
            else:
                return {"success": False, "error": "Failed to write file", "file_path": file_path}
                
        except ImportError:
            logger.error("python-docx is required for DOCX operations. Install with: pip install python-docx")
            return {"success": False, "error": "python-docx not available", "file_path": file_path}
        except Exception as e:
            logger.error(f"Error saving DOCX file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    def _append_docx(self, file_path: str, content: Any, **kwargs) -> Dict[str, Any]:
        """
        Append content to an existing DOCX file.
        
        Input: String → paragraph, List → paragraphs/bullet points
        Output: Appends to existing DOCX document (creates new if doesn't exist)
        """
        try:
            from docx import Document
            from io import BytesIO
            
            # Read existing document if it exists
            if self._exists_raw(file_path):
                existing_bytes = self._read_raw(file_path, **kwargs)
                doc = Document(BytesIO(existing_bytes))
            else:
                # Create new document if file doesn't exist
                doc = Document()
            
            # Handle different content types for appending
            if isinstance(content, str):
                # Add single paragraph
                doc.add_paragraph(content)
            
            elif isinstance(content, list):
                # List of paragraphs or lists
                for item in content:
                    if isinstance(item, str):
                        # Add as plain text paragraph
                        doc.add_paragraph(item)
                    elif isinstance(item, list):
                        # Add as bullet list
                        for list_item in item:
                            if isinstance(list_item, str):
                                # Try to use bullet style, fall back to normal paragraph with bullet
                                try:
                                    doc.add_paragraph(list_item, style='List Bullet')
                                except Exception:
                                    doc.add_paragraph(f"• {list_item}")
            
            # Save document to bytes
            doc_buffer = BytesIO()
            doc.save(doc_buffer)
            doc_bytes = doc_buffer.getvalue()
            
            # Use raw write method
            success = self._write_raw(file_path, doc_bytes, **kwargs)
            
            if success:
                res = {
                    "success": True,
                    "message": "DOCX file appended successfully",
                    "file_path": file_path
                }
                if self.return_file_url:
                    res["url"] = self._get_file_url(file_path)
                return res
            else:
                return {"success": False, "error": "Failed to write file", "file_path": file_path}
                
        except ImportError:
            logger.error("python-docx is required for DOCX operations. Install with: pip install python-docx")
            return {"success": False, "error": "python-docx not available", "file_path": file_path}
        except Exception as e:
            logger.error(f"Error appending to DOCX file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    # Removed legacy DOC handlers as per project decision to support DOCX only

    # Pickle file handlers
    def _save_pickle(self, file_path: str, content: Any, **kwargs) -> Dict[str, Any]:
        """Save pickle content to a file"""
        try:
            # Convert content to bytes using pickle
            content_bytes = pickle.dumps(content)
            
            # Use raw write method
            success = self._write_raw(file_path, content_bytes, **kwargs)
            
            if success:
                res = {
                    "success": True,
                    "message": "Pickle file saved successfully",
                    "file_path": file_path
                }
                if self.return_file_url:
                    res["url"] = self._get_file_url(file_path)
                return res
            else:
                return {"success": False, "error": "Failed to write file", "file_path": file_path}
        except Exception as e:
            logger.error(f"Error saving pickle file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    def _read_pickle(self, file_path: Optional[str] = None, content: Any = None, **kwargs) -> Dict[str, Any]:
        """Read pickle content from a file"""
        try:
            # Use raw read method
            if content is None:
                if not file_path:
                    return {"success": False, "error": "Either file_path or content must be provided"}
                content = self._read_raw(file_path, **kwargs)
            
            if not isinstance(content, bytes):
                return {"success": False, "error": "Content must be bytes"}

            content_bytes = content
            
            # Parse pickle content
            content = pickle.loads(content_bytes)
            
            base = {
                "success": True,
                "content": content,
                "file_path": file_path
            }
            return base
        except Exception as e:
            logger.error(f"Error reading pickle file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    def _append_pickle(self, file_path: str, content: Any, **kwargs) -> Dict[str, Any]:
        """Append content to pickle file (for lists)"""
        try:
            # Read existing content
            existing_content = []
            if self._exists_raw(file_path):
                existing_bytes = self._read_raw(file_path, **kwargs)
                existing_content = pickle.loads(existing_bytes)
            
            # Merge content
            if isinstance(existing_content, list):
                if isinstance(content, list):
                    existing_content.extend(content)
                else:
                    existing_content.append(content)
            elif isinstance(existing_content, dict):
                if isinstance(content, dict):
                    existing_content.update(content)
                elif isinstance(content, list):
                    existing_content["appended_list"] = content
                else:
                    existing_content["appended_value"] = content
            else:
                existing_content = [existing_content]
                if isinstance(content, list):
                    existing_content.extend(content)
                else:
                    existing_content.append(content)
            
            # Convert to bytes and write
            content_bytes = pickle.dumps(existing_content)
            success = self._write_raw(file_path, content_bytes, **kwargs)
            
            if success:
                res = {
                    "success": True,
                    "message": f"Content appended to pickle file {file_path}",
                    "file_path": file_path
                }
                if self.return_file_url:
                    res["url"] = self._get_file_url(file_path)
                return res
            else:
                return {"success": False, "error": "Failed to append to file", "file_path": file_path}
        except Exception as e:
            logger.error(f"Error appending to pickle file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    # PDF file handlers
    def _save_pdf(self, file_path: str, content: Union[str, bytes], **kwargs) -> Dict[str, Any]:
        """Save content to a PDF file"""
            
        try:
            if isinstance(content, bytes):
                content_bytes = content
            elif isinstance(content, str):
                try:
                    from markdown_pdf import MarkdownPdf, Section
                    import tempfile
                    pdf = MarkdownPdf(toc_level=2)
                    pdf.add_section(Section(content))
                    
                    tmp_path = None
                    try:
                        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
                            tmp_path = tmp.name
                    
                        pdf.save(tmp_path)
                    
                        with open(tmp_path, "rb") as f:
                            content_bytes = f.read()
                    finally:
                        if tmp_path and os.path.exists(tmp_path):
                            os.unlink(tmp_path)
                    
                except ImportError:
                    return {"success": False, "error": "markdown-pdf library not available for PDF creation", "file_path": file_path}
            else:
                content_bytes = str(content).encode("utf-8")

            success = self._write_raw(file_path, content_bytes, **kwargs)
            if success:
                res = {
                    "success": True,
                    "message": "PDF file saved",
                    "file_path": file_path
                }
                if self.return_file_url:
                    res["url"] = self._get_file_url(file_path)
                return res
            else:
                return {"success": False, "error": "Failed to write file", "file_path": file_path}
        except Exception as e:
            logger.error(f"Error saving PDF file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    def _read_pdf(self, file_path: Optional[str] = None, content: Any = None, **kwargs) -> Dict[str, Any]:
        """Read content from a PDF file"""
        if not PDF_AVAILABLE:
            return {"success": False, "error": "PyMuPDF library not available", "file_path": file_path}
        try:
            if content is None:
                if not file_path:
                    return {"success": False, "error": "Either file_path or content must be provided"}
                content = self._read_raw(file_path, **kwargs)
            
            if not isinstance(content, bytes):
                return {"success": False, "error": "Content must be bytes"}

            content_bytes = content
            doc = pymupdf.open(stream=content_bytes, filetype="pdf")
            all_text = []
            for page in doc:
                all_text.append(page.get_text())
            doc.close()
            text = "\n\n".join(all_text)
            base = {
                "success": True,
                "content": text,
                "file_path": file_path
            }
            return base
        except Exception as e:
            logger.error(f"Error reading PDF file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    # PowerPoint file handlers
    def _read_pptx(self, file_path: Optional[str] = None, content: Any = None, **kwargs) -> Dict[str, Any]:
        """Read content from a PowerPoint presentation"""
        if not PPTX_AVAILABLE:
            return {"success": False, "error": "python-pptx library not available", "file_path": file_path}
        
        try:
            from io import BytesIO
            
            if content is None:
                if not file_path:
                    return {"success": False, "error": "Either file_path or content must be provided"}
                content = self._read_raw(file_path, **kwargs)
            
            if not isinstance(content, bytes):
                return {"success": False, "error": "Content must be bytes"}

            content_bytes = content
            prs = Presentation(BytesIO(content_bytes))
            
            slides_content = []
            for i, slide in enumerate(prs.slides):
                slide_data = {
                    "slide_number": i + 1,
                    "title": "",
                    "content": []
                }
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Check if this is likely a title (first text shape or larger font)
                        if not slide_data["title"] and (i == 0 or len(shape.text) < 100):
                            slide_data["title"] = shape.text.strip()
                        else:
                            slide_data["content"].append(shape.text.strip())
                
                slides_content.append(slide_data)
            
            # Create a text summary of all slides
            text_content = []
            for slide in slides_content:
                if slide["title"]:
                    text_content.append(f"Slide {slide['slide_number']}: {slide['title']}")
                else:
                    text_content.append(f"Slide {slide['slide_number']}:")
                
                for content in slide["content"]:
                    text_content.append(f"  {content}")
                text_content.append("")  # Empty line between slides
            
            base = {
                "success": True,
                "content": "\n".join(text_content),
                "slides": slides_content,
                "slide_count": len(slides_content),
                "file_path": file_path
            }
            return base
            
        except Exception as e:
            logger.error(f"Error reading PPTX file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    def _save_pptx(self, file_path: str, content: Any, **kwargs) -> Dict[str, Any]:
        """
        Create a new PowerPoint presentation.
        
        Input: String → single slide with content, List → multiple slides/bullet points
        Output: Creates new PPTX document
        """
        if not PPTX_AVAILABLE:
            return {"success": False, "error": "python-pptx library not available", "file_path": file_path}
        
        try:
            from io import BytesIO
            
            prs = Presentation()
            
            if isinstance(content, str):
                # Single slide with content
                slide_layout = prs.slide_layouts[1]  # Title and Content layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Add title
                title = slide.shapes.title
                title.text = "Slide 1"
                
                # Add content
                content_placeholder = slide.placeholders[1]
                content_placeholder.text = content
                
            elif isinstance(content, list):
                for i, item in enumerate(content):
                    slide_layout = prs.slide_layouts[1]  # Title and Content layout
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # Add title
                    title = slide.shapes.title
                    title.text = f"Slide {i + 1}"
                    
                    if isinstance(item, str):
                        # Single content item
                        content_placeholder = slide.placeholders[1]
                        content_placeholder.text = item
                    elif isinstance(item, list):
                        # Bullet points
                        content_placeholder = slide.placeholders[1]
                        text_frame = content_placeholder.text_frame
                        text_frame.clear()
                        
                        for j, bullet in enumerate(item):
                            p = text_frame.paragraphs[0] if j == 0 else text_frame.add_paragraph()
                            p.text = str(bullet)
                            p.level = 0
                    elif isinstance(item, dict) and "title" in item:
                        # Slide with custom title and content
                        title.text = item["title"]
                        slide_content = item.get("content", "")
                        
                        if isinstance(slide_content, str):
                            content_placeholder = slide.placeholders[1]
                            content_placeholder.text = slide_content
                        elif isinstance(slide_content, list):
                            content_placeholder = slide.placeholders[1]
                            text_frame = content_placeholder.text_frame
                            text_frame.clear()
                            
                            for j, bullet in enumerate(slide_content):
                                p = text_frame.paragraphs[0] if j == 0 else text_frame.add_paragraph()
                                p.text = str(bullet)
                                p.level = 0
            else:
                # Fallback: create single slide with string representation
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = "Content"
                content_placeholder = slide.placeholders[1]
                content_placeholder.text = str(content)
            
            # Save to BytesIO
            buffer = BytesIO()
            prs.save(buffer)
            content_bytes = buffer.getvalue()
            
            success = self._write_raw(file_path, content_bytes, **kwargs)
            if success:
                res = {
                    "success": True,
                    "message": f"PPTX file saved with {len(prs.slides)} slides",
                    "file_path": file_path,
                    "slide_count": len(prs.slides)
                }
                if self.return_file_url:
                    res["url"] = self._get_file_url(file_path)
                return res
            else:
                return {"success": False, "error": "Failed to write file", "file_path": file_path}
                
        except Exception as e:
            logger.error(f"Error saving PPTX file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    def _append_pptx(self, file_path: str, content: Any, **kwargs) -> Dict[str, Any]:
        """
        Append slides to an existing PowerPoint presentation.
        
        Input: String → single slide with content, List → multiple slides/bullet points
        Output: Appends to existing PPTX document (creates new if doesn't exist)
        """
        if not PPTX_AVAILABLE:
            return {"success": False, "error": "python-pptx library not available", "file_path": file_path}
        
        try:
            from io import BytesIO
            
            # Try to read existing presentation
            try:
                content_bytes = self._read_raw(file_path, **kwargs)
                prs = Presentation(BytesIO(content_bytes))
                initial_slide_count = len(prs.slides)
            except Exception:
                # File doesn't exist or can't be read, create new presentation
                prs = Presentation()
                initial_slide_count = 0
            
            # Add new content
            if isinstance(content, str):
                # Single slide with content
                slide_layout = prs.slide_layouts[1]  # Title and Content layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Add title
                title = slide.shapes.title
                title.text = f"Slide {len(prs.slides)}"
                
                # Add content
                content_placeholder = slide.placeholders[1]
                content_placeholder.text = content
                
            elif isinstance(content, list):
                for i, item in enumerate(content):
                    slide_layout = prs.slide_layouts[1]  # Title and Content layout
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # Add title
                    title = slide.shapes.title
                    title.text = f"Slide {len(prs.slides)}"
                    
                    if isinstance(item, str):
                        # Single content item
                        content_placeholder = slide.placeholders[1]
                        content_placeholder.text = item
                    elif isinstance(item, list):
                        # Bullet points
                        content_placeholder = slide.placeholders[1]
                        text_frame = content_placeholder.text_frame
                        text_frame.clear()
                        
                        for j, bullet in enumerate(item):
                            p = text_frame.paragraphs[0] if j == 0 else text_frame.add_paragraph()
                            p.text = str(bullet)
                            p.level = 0
                    elif isinstance(item, dict) and "title" in item:
                        # Slide with custom title and content
                        title.text = item["title"]
                        slide_content = item.get("content", "")
                        
                        if isinstance(slide_content, str):
                            content_placeholder = slide.placeholders[1]
                            content_placeholder.text = slide_content
                        elif isinstance(slide_content, list):
                            content_placeholder = slide.placeholders[1]
                            text_frame = content_placeholder.text_frame
                            text_frame.clear()
                            
                            for j, bullet in enumerate(slide_content):
                                p = text_frame.paragraphs[0] if j == 0 else text_frame.add_paragraph()
                                p.text = str(bullet)
                                p.level = 0
            else:
                # Fallback: create single slide with string representation
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = f"Slide {len(prs.slides)}"
                content_placeholder = slide.placeholders[1]
                content_placeholder.text = str(content)
            
            # Save to BytesIO
            buffer = BytesIO()
            prs.save(buffer)
            content_bytes = buffer.getvalue()
            
            success = self._write_raw(file_path, content_bytes, **kwargs)
            if success:
                new_slides = len(prs.slides) - initial_slide_count
                res = {
                    "success": True,
                    "message": f"Added {new_slides} slides to PPTX file (total: {len(prs.slides)} slides)",
                    "file_path": file_path,
                    "slides_added": new_slides,
                    "total_slides": len(prs.slides)
                }
                if self.return_file_url:
                    res["url"] = self._get_file_url(file_path)
                return res
            else:
                return {"success": False, "error": "Failed to write file", "file_path": file_path}
                
        except Exception as e:
            logger.error(f"Error appending to PPTX file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    # Image file handlers
    def _save_image(self, file_path: str, content: Any, **kwargs) -> Dict[str, Any]:
        """Save image content to a file"""
        try:
            from io import BytesIO
            
            # Check if content is already a PIL Image object
            if hasattr(content, 'save') and callable(getattr(content, 'save', None)):
                # Content is a PIL Image object - save to BytesIO first
                buffer = BytesIO()
                content.save(buffer, format=content.format or 'PNG')
                content_bytes = buffer.getvalue()
                
                # Use raw write method
                success = self._write_raw(file_path, content_bytes, **kwargs)
                
                if success:
                    res = {
                        "success": True,
                        "message": "Image saved successfullly",
                        "file_path": file_path,
                        "format": content.format,
                        "size": content.size
                    }
                    if self.return_file_url:
                        res["url"] = self._get_file_url(file_path)
                    return res
                else:
                    return {"success": False, "error": "Failed to write file", "file_path": file_path}
            elif isinstance(content, bytes):
                # Content is binary image data
                success = self._write_raw(file_path, content, **kwargs)
                
                if success:
                    res = {
                        "success": True,
                        "message": "Image saved successfullly",
                        "file_path": file_path
                    }
                    if self.return_file_url:
                        res["url"] = self._get_file_url(file_path)
                    return res
                else:
                    return {"success": False, "error": "Failed to write file", "file_path": file_path}
            elif isinstance(content, str) and Path(content).exists():
                # Content is a file path to an existing image - read and write
                with open(content, 'rb') as f:
                    content_bytes = f.read()
                
                success = self._write_raw(file_path, content_bytes, **kwargs)
                
                if success:
                    res = {
                        "success": True,
                        "message": f"Image copied from {content} to {file_path}",
                        "file_path": file_path
                    }
                    if self.return_file_url:
                        res["url"] = self._get_file_url(file_path)
                    return res
                else:
                    return {"success": False, "error": "Failed to write file", "file_path": file_path}
            else:
                return {"success": False, "error": "Content must be a PIL Image object, binary data, or valid file path"}
                
        except Exception as e:
            logger.error(f"Error saving image file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    def _read_image(self, file_path: Optional[str] = None, content: Any = None, **kwargs) -> Dict[str, Any]:
        """Read image and return PIL Image object"""
        if not PILLOW_AVAILABLE:
            return {"success": False, "error": "Pillow library not available"}
        
        try:
            from io import BytesIO
            
            # Use raw read method
            if content is None:
                if not file_path:
                    return {"success": False, "error": "Either file_path or content must be provided"}
                content = self._read_raw(file_path, **kwargs)
            
            if not isinstance(content, bytes):
                return {"success": False, "error": "Content must be bytes"}

            content_bytes = content
            
            # Open image from bytes
            with Image.open(BytesIO(content_bytes)) as img:
                # Convert to RGB if necessary for consistency
                if img.mode in ('RGBA', 'LA', 'P'):
                    img = img.convert('RGB')
                
                metadata = {
                    "format": img.format,
                    "mode": img.mode,
                    "size": img.size,
                    "width": img.width,
                    "height": img.height
                }
                
                base = {
                    "success": True,
                    "content": img,  # Return the PIL Image object
                    "metadata": metadata,
                    "file_path": file_path
                }
            return base
                
        except Exception as e:
            logger.error(f"Error reading image file {file_path}: {str(e)}")
            return {"success": False, "error": str(e), "file_path": file_path}
    
    # Download methods
    def _extract_filename_from_url(self, url: str, content_type: str = None) -> str:
        """Extract filename from URL or generate one based on content type."""
        parsed_url = urlparse(url)
        
        # Try to get filename from URL path
        path = unquote(parsed_url.path)
        if path and '/' in path:
            filename = path.split('/')[-1]
            if filename and '.' in filename:
                return filename
        
        # Try to get filename from query parameters
        query_params = parse_qs(parsed_url.query)
        for param in ['filename', 'file', 'name']:
            if param in query_params and query_params[param]:
                filename = query_params[param][0]
                if filename:
                    return filename
        
        # Generate filename based on content type
        if content_type:
            extension = mimetypes.guess_extension(content_type.split(';')[0])
            if extension:
                timestamp = int(time.time())
                return f"download_{timestamp}{extension}"
        
        # Default filename
        timestamp = int(time.time())
        return f"download_{timestamp}"

    def _get_content_type_from_response(self, response) -> str:
        """Extract content type from response headers."""
        content_type = response.headers.get('content-type', '')
        return content_type.split(';')[0].strip()

    def _sanitize_filename(self, filename: str) -> str:
        """Sanitize filename to remove invalid characters."""
        # Remove or replace invalid characters
        filename = re.sub(r'[<>:"/\\|?*]', '_', filename)
        # Remove leading/trailing spaces and dots
        filename = filename.strip(' .')
        # Ensure filename is not empty
        if not filename:
            filename = f"download_{int(time.time())}"
        return filename

    def _download_content(self, url: str, filename: str = None, timeout: int = 30, 
                         max_retries: int = 3, headers: Dict[str, str] = None, contain_file_path=True, memory_only: bool = False, **kwargs) -> Dict[str, Any]:  
        """
        Download content from URL and save it using the save method.
        
        Args:
            url: URL to download from
            filename: Optional pure filename (without extension) to save as
            timeout: Request timeout in seconds
            max_retries: Maximum number of retry attempts
            headers: Optional headers to include in request
            memory_only: If True, return content bytes directly without saving to storage
            
        Returns:
            Dict containing success status, file path (or content if memory_only), and metadata
        """
        try:
            # Prepare headers
            request_headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
            }
            if headers:
                request_headers.update(headers)
            
            # Download with retries
            last_exception = None
            for attempt in range(max_retries):
                try:
                    logger.info(f"Downloading from {url} (attempt {attempt + 1}/{max_retries})")
                    
                    response = requests.get(url, headers=request_headers, timeout=timeout, stream=True)
                    response.raise_for_status()
                    
                    # Get content type
                    content_type = self._get_content_type_from_response(response)
                    
                    # Download content first
                    content = b''
                    # total_size = int(response.headers.get('content-length', 0))
                    downloaded_size = 0
                    
                    for chunk in response.iter_content(chunk_size=8192):
                        if chunk:
                            content += chunk
                            downloaded_size += len(chunk)
                    
                    # Determine final filename with robust extension resolution
                    head = content[:65536]
                    inferred_ext = self._resolve_extension(response, url, filename, head)
                    stem = self._sanitize_filename(filename) if filename else self._sanitize_filename(self._extract_filename_from_url(url, content_type).split('.')[0])
                    final_filename = stem + inferred_ext
                    
                    # Sanitize final filename
                    final_filename = self._sanitize_filename(final_filename)
                    
                    if memory_only:
                        logger.info("Successfully downloaded to memory")
                        return {
                            'success': True,
                            'url': url,
                            'content': content,
                            'content_type': content_type,
                            'size': len(content),
                            'filename': final_filename,
                            'extension': inferred_ext
                        }

                    # Save using the save method which handles bytes properly
                    save_result = self.save(final_filename, content, return_path=True, **kwargs)
                    
                    if save_result.get('success', False):
                        logger.info("Successfully downloaded")
                        result = {
                            'success': True,
                            'url': save_result.get("url", None),
                            'content_type': content_type,
                            'size': len(content),
                            'filename': final_filename
                        }
                        if contain_file_path and "file_path" in save_result:
                            result["file_path"] = self.translate_out(save_result["file_path"])
                        return result
                    else:
                        return {
                            'success': False,
                            'error': f"Failed to save downloaded content: {save_result.get('error', 'Unknown error')}",
                            'url': url
                        }
                        
                except requests.exceptions.RequestException as e:
                    last_exception = e
                    if attempt < max_retries - 1:
                        wait_time = 2 ** attempt  # Exponential backoff
                        logger.warning(f"Download attempt {attempt + 1} failed: {str(e)}. Retrying in {wait_time} seconds...")
                        time.sleep(wait_time)
                    continue
                    
            # All retries failed
            error_msg = f"Failed to download after {max_retries} attempts. Last error: {str(last_exception)}"
            logger.error(error_msg)
            return {
                'success': False,
                'error': error_msg,
                'url': url
            }
            
        except Exception as e:
            import traceback
            error_msg = f"Unexpected error during download: {str(e)}\n{traceback.format_exc()}"
            logger.error(error_msg)
            return {
                'success': False,
                'error': error_msg,
                'url': url
            }

    # Placeholder for future database integration
    def _get_database_connection(self, db_type: str, connection_string: str) -> Any:
        """Placeholder for future database integration"""
        # This will be implemented when adding database support
        raise NotImplementedError("Database integration not yet implemented") 
    
    def _get_file_url(self, file_path: str) -> str:
        """Create URL for others access"""
        raise NotImplementedError("URL convertion methods not implemented")

    # ____________________ DOWNLOAD HELPERS ____________________ #
    def _parse_content_disposition_filename(self, response) -> Optional[str]:
        cd = response.headers.get('content-disposition') or response.headers.get('Content-Disposition')
        if not cd:
            return None
        try:
            # Basic parsing for filename and RFC 5987 filename*
            # Examples:
            # Content-Disposition: attachment; filename="paper.pdf"
            # Content-Disposition: attachment; filename*=UTF-8''paper.pdf
            parts = [p.strip() for p in cd.split(';')]
            kv = {}
            for p in parts[1:]:
                if '=' in p:
                    k, v = p.split('=', 1)
                    kv[k.strip().lower()] = v.strip().strip('"')
            if 'filename*' in kv:
                val = kv['filename*']
                # RFC 5987: charset''encoded_filename
                if "''" in val:
                    _, encoded = val.split("''", 1)
                    from urllib.parse import unquote
                    return self._sanitize_filename(unquote(encoded))
                return self._sanitize_filename(val)
            if 'filename' in kv:
                return self._sanitize_filename(kv['filename'])
        except Exception:
            return None
        return None

    def _guess_extension_from_head(self, head: bytes) -> Optional[str]:
        try:
            if head.startswith(b'%PDF-'):
                return '.pdf'
            if head.startswith(b'\x89PNG\r\n\x1a\n'):
                return '.png'
            if head.startswith(b'\xFF\xD8\xFF'):
                return '.jpg'
            if head.startswith(b'GIF87a') or head.startswith(b'GIF89a'):
                return '.gif'
            if head[:4] == b'RIFF' and b'WEBP' in head[:16]:
                return '.webp'
            if head[:2] == b'PK':
                # ZIP-based container — try to identify OOXML types by peeking
                # at the ZIP central directory (uses only the bytes already in memory).
                import zipfile
                import io
                try:
                    with zipfile.ZipFile(io.BytesIO(head)) as zf:
                        names = zf.namelist()
                        if any(n.startswith('xl/') for n in names):
                            return '.xlsx'
                        if any(n.startswith('word/') for n in names):
                            return '.docx'
                        if any(n.startswith('ppt/') for n in names):
                            return '.pptx'
                except Exception:
                    pass
                return '.zip'
        except Exception:
            pass
        return None

    def _resolve_extension(self, response, url: str, user_stem: Optional[str], head: bytes) -> str:
        # 1) Content-Disposition filename
        cd_name = self._parse_content_disposition_filename(response)
        if cd_name:
            _, ext = os.path.splitext(cd_name)
            if ext and ext.lower() in mimetypes.types_map:
                return ext.lower()
        # 2) URL extension (final URL after redirects) — check before magic sniff
        #    so that an explicit extension in the URL (e.g. .xlsx) is not overridden
        #    by an ambiguous magic-byte match (e.g. PK → .zip for OOXML files).
        #    Pass content_type=None so _extract_filename_from_url only reads the
        #    actual URL path and never falls back to generating a CT-derived extension
        #    (e.g. application/octet-stream → .bin) which would incorrectly skip
        #    the magic sniff for extension-less URLs.
        ct = self._get_content_type_from_response(response)
        temp = self._extract_filename_from_url(response.url if getattr(response, 'url', None) else url, None)
        url_ext = os.path.splitext(temp)[1].lower()
        if url_ext and url_ext in mimetypes.types_map:
            return url_ext
        # 3) Magic sniff
        magic_ext = self._guess_extension_from_head(head or b'')
        if magic_ext:
            return magic_ext
        # 4) Content-Type mapping
        ct_ext = mimetypes.guess_extension(ct) or None
        if ct_ext:
            return ct_ext
        # 5) Fallback
        return '.bin'
